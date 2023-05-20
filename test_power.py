import collections.abc
from pptx import Presentation

prs = Presentation()
title_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

subtitle.text = "Hello, World!"
subtitle.text = "python-pptx was here!"

def pptx2text(file_path):
    texts = [] # 抽出したテキストデータを格納する空リスト
    prs = pptx.Presentation(file_path)
    # スライドごとにテキストデータを抽出する
    for sld in prs.slides:
        for shape in sld.shapes:
            # shapeに含まれるテキストデータを抽出
            if shape.has_text_frame:
                for text in shape.text.splitlines():
                    texts.append(text)
            # tableに含まれるテキストデータを抽出
            if shape.has_table: 
                for cell in shape.table.iter_cells():
                    for text in cell.text.splitlines():
                        texts.append(text)
    return texts

print(len('test.pptx'))
prs.save('test2.pptx')